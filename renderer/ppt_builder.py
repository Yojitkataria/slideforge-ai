"""PPT builder for SlideForge AI slide plans."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from renderer.layout_engine import get_layout, split_content

TEMPLATE_PATH = Path(__file__).resolve().parents[1] / "templates" / "template.pptx"
DEFAULT_TEMPLATE_PATH = Path(__file__).resolve().parents[1] / "templates" / "default.pptx"

# Mapped from detected template layouts:
# 0 Cover, 1 Divider, 2 Blank, 3 Title only, 4 1_Thank you
LAYOUT_MAP = {
    "title": 0,
    "agenda": 1,
    "summary": 1,
    "content": 1,
    "chart": 1,
    "table": 1,
    "comparison": 1,
    "process": 1,
    "timeline": 1,
    "two_column": 1,
    "conclusion": 4,
}


def build_ppt(slide_plan: dict[str, Any], output_file: str, template: str | None = None) -> str:
    """
    Convert slide plan JSON into a PPTX file.

    Args:
        slide_plan: Dict containing {"slides": [...]}
        output_file: Target pptx file path

    Returns:
        Output file path
    """
    if template:
        prs = Presentation(str(template))
    else:
        fallback = DEFAULT_TEMPLATE_PATH if DEFAULT_TEMPLATE_PATH.exists() else TEMPLATE_PATH
        prs = Presentation(str(fallback))
    slides = slide_plan.get("slides", [])

    for slide_spec in slides:
        slide_type = str(slide_spec.get("type", "content")).lower()
        layout_config = get_layout(slide_type)

        if slide_type in {"title", "agenda", "summary", "content", "conclusion"}:
            _add_text_slides(prs, slide_spec, layout_config)
        elif slide_type == "chart":
            layout_index = _resolve_layout_index(prs, slide_type)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            title = _safe_text(slide_spec.get("title")) or "Chart"
            _set_title_text(slide.shapes.title or _get_first_text_placeholder(slide), title, layout_config)
            add_chart(slide, slide_spec.get("data", []), layout_config)
        elif slide_type == "table":
            layout_index = _resolve_layout_index(prs, slide_type)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            title = _safe_text(slide_spec.get("title")) or "Table"
            _set_title_text(slide.shapes.title or _get_first_text_placeholder(slide), title, layout_config)
            add_table(slide, slide_spec.get("table", []), layout_config)
        elif slide_type in {"comparison", "two_column", "process", "timeline"}:
            layout_index = _resolve_layout_index(prs, slide_type)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            title = _safe_text(slide_spec.get("title")) or "Visual"
            _set_title_text(slide.shapes.title or _get_first_text_placeholder(slide), title, layout_config)
            content = _to_string_list(slide_spec.get("content", []))
            if slide_type in {"comparison", "two_column"}:
                add_comparison(slide, content)
            elif slide_type == "process":
                add_process(slide, content)
            elif slide_type == "timeline":
                add_timeline(slide, content)
        else:
            _add_text_slides(prs, slide_spec, layout_config)

    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)


def add_chart(slide: Any, data: list[dict[str, Any]], layout: dict[str, Any]) -> None:
    """Create a bar chart under title using template content area."""
    categories: list[str] = []
    values: list[float] = []
    for item in data:
        if not isinstance(item, dict):
            continue
        label = _safe_text(item.get("label"))
        value = _to_float(item.get("value"))
        if not label or value is None:
            continue
        categories.append(label)
        values.append(value)

    body_shape = _get_body_placeholder(slide)
    if not categories:
        if body_shape is not None:
            _fill_body_bullets(body_shape, ["No chart data available."], layout)
        return
    if body_shape is None:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series 1", values)

    x, y, cx, cy = body_shape.left, body_shape.top, body_shape.width, body_shape.height
    body_shape._element.getparent().remove(body_shape._element)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True


def add_table(slide: Any, table_data: list[dict[str, Any]], layout: dict[str, Any]) -> None:
    """Create table in template content area."""
    if not isinstance(table_data, list) or not table_data:
        body_shape = _get_body_placeholder(slide)
        if body_shape is not None:
            _fill_body_bullets(body_shape, ["No table data available."], layout)
        return

    headers = _extract_table_headers(table_data)
    if not headers:
        body_shape = _get_body_placeholder(slide)
        if body_shape is not None:
            _fill_body_bullets(body_shape, ["Table data format is invalid."], layout)
        return

    body_shape = _get_body_placeholder(slide)
    if body_shape is None:
        return

    rows = len(table_data) + 1
    cols = len(headers)
    x, y, cx, cy = body_shape.left, body_shape.top, body_shape.width, body_shape.height
    body_shape._element.getparent().remove(body_shape._element)
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table

    for col, header in enumerate(headers):
        table.cell(0, col).text = header

    for row, item in enumerate(table_data, start=1):
        for col, key in enumerate(headers):
            table.cell(row, col).text = _safe_text(item.get(key))


def add_comparison(slide: Any, content: list[str]) -> None:
    left, top, width, height = _get_content_region(
        slide,
        fallback=(Inches(0.5), Inches(2.0), Inches(9.0), Inches(3.2)),
    )
    gap = Inches(0.3)
    col_width = max(Inches(1.5), int((width - gap) / 2))
    right = left + col_width + gap

    box1 = slide.shapes.add_textbox(left, top, col_width, height)
    box2 = slide.shapes.add_textbox(right, top, col_width, height)

    tf1 = box1.text_frame
    tf2 = box2.text_frame
    tf1.clear()
    tf2.clear()

    mid = max(1, len(content) // 2) if content else 0
    left_items = content[:mid] if content else ["Left"]
    right_items = content[mid:] if content and mid < len(content) else ["Right"]

    for idx, item in enumerate(left_items):
        p = tf1.paragraphs[0] if idx == 0 else tf1.add_paragraph()
        p.text = item
        p.space_after = Pt(8)

    for idx, item in enumerate(right_items):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        p.text = item
        p.space_after = Pt(8)


def add_process(slide: Any, content: list[str]) -> None:
    left, top, width, height = _get_content_region(
        slide,
        fallback=(Inches(0.5), Inches(2.0), Inches(9.0), Inches(3.2)),
    )
    steps = content[:5] if content else ["Step 1", "Step 2", "Step 3"]
    _add_horizontal_flow_boxes(slide, steps, left, top, width, height)


def add_timeline(slide: Any, content: list[str]) -> None:
    left, top, width, height = _get_content_region(
        slide,
        fallback=(Inches(0.5), Inches(2.2), Inches(9.0), Inches(2.8)),
    )
    items = content[:5] if content else ["Phase 1", "Phase 2", "Phase 3"]
    _add_horizontal_flow_boxes(slide, items, left, top, width, height)


def _add_text_slides(prs: Presentation, slide_spec: dict[str, Any], layout: dict[str, Any]) -> None:
    slide_type = str(slide_spec.get("type", "content")).lower()
    title = _safe_text(slide_spec.get("title")) or "Slide"
    content_items = _to_string_list(slide_spec.get("content", []))
    max_bullets = int(layout.get("rules", {}).get("max_bullets_per_slide", 5))
    chunks = split_content(content_items, max_bullets) or [[" "]]

    for idx, chunk in enumerate(chunks, start=1):
        layout_index = _resolve_layout_index(prs, slide_type)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        chunk_title = title if idx == 1 else f"{title} (Part {idx})"
        title_shape = slide.shapes.title or _get_first_text_placeholder(slide)
        body_shape = _get_body_placeholder(slide)

        _set_title_text(title_shape, chunk_title, layout)
        if body_shape is not None:
            _fill_body_bullets(body_shape, chunk, layout)


def _extract_table_headers(rows: list[Any]) -> list[str]:
    for row in rows:
        if isinstance(row, dict) and row:
            return [str(key) for key in row.keys()]
    return []


def _set_title_text(title_shape: Any, text: str, layout: dict[str, Any]) -> None:
    if title_shape is None or not hasattr(title_shape, "text_frame"):
        return
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(float(layout["title"]["font_size"]))
    p.font.bold = True
    p.alignment = _to_alignment(str(layout["title"].get("align", "center")))


def _fill_body_bullets(body_shape: Any, bullets: list[str], layout: dict[str, Any]) -> None:
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    spacing_pt = Pt(float(layout["content"]["spacing"]) * 72.0)
    body_size = Pt(float(layout["content"]["font_size"]))
    align = _to_alignment(str(layout["content"].get("align", "left")))

    content = bullets or [" "]
    for idx, bullet in enumerate(content):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = body_size
        p.alignment = align
        p.space_after = spacing_pt


def _to_string_list(value: Any) -> list[str]:
    if not isinstance(value, list):
        return []
    return [str(item) for item in value if str(item).strip()]


def _safe_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _to_float(value: Any) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _get_body_placeholder(slide: Any) -> Any:
    # Equivalent to "placeholders[1]" by position, not placeholder idx key.
    placeholders = list(slide.placeholders)
    if len(placeholders) > 1:
        return placeholders[1]
    if len(placeholders) > 0:
        return placeholders[0]
    return None


def _get_first_text_placeholder(slide: Any) -> Any:
    for ph in slide.placeholders:
        if hasattr(ph, "text_frame"):
            return ph
    return None


def _get_content_region(
    slide: Any,
    fallback: tuple[int, int, int, int],
) -> tuple[int, int, int, int]:
    body_shape = _get_body_placeholder(slide)
    if body_shape is None:
        return fallback
    left, top, width, height = body_shape.left, body_shape.top, body_shape.width, body_shape.height
    body_shape._element.getparent().remove(body_shape._element)
    return left, top, width, height


def _add_horizontal_flow_boxes(
    slide: Any,
    items: list[str],
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    if not items:
        return

    count = len(items)
    columns = min(4, count)
    rows = (count + columns - 1) // columns

    col_gap = Inches(0.18)
    row_gap = Inches(0.2)
    box_width = max(Inches(1.3), int((width - col_gap * (columns - 1)) / columns))
    box_height = max(Inches(0.7), int((height - row_gap * (rows - 1)) / rows))

    for idx, item in enumerate(items):
        row = idx // columns
        col = idx % columns
        x = left + col * (box_width + col_gap)
        y = top + row * (box_height + row_gap)
        box = slide.shapes.add_textbox(x, y, box_width, box_height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)


def _resolve_layout_index(prs: Presentation, slide_type: str) -> int:
    mapped = LAYOUT_MAP.get(slide_type, LAYOUT_MAP["content"])
    if 0 <= mapped < len(prs.slide_layouts):
        return mapped
    return 0


def _to_alignment(value: str) -> Any:
    normalized = str(value).strip().lower()
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT
